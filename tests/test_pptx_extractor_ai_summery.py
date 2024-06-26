import os
import sys
import pytest
import json
from pptx import Presentation

# הוספת המסלול לתיקיית החבילה לספריית ה-PATH
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../explainer/scripts')))

from explainer.scripts.pptx_extractor import extract_text_from_presentation, text_to_json_file
from explainer.scripts.async_tasks import process_presentation
from explainer.scripts.ai_api import load_api_key

OUTPUTS_FOLDER = os.path.join(os.path.dirname(__file__), 'outputs')
DEMO_FILES_FOLDER = os.path.join(os.path.dirname(__file__), 'demo_files')
os.makedirs(OUTPUTS_FOLDER, exist_ok=True)
os.makedirs(DEMO_FILES_FOLDER, exist_ok=True)

@pytest.mark.asyncio
async def test_presentation_summary():
    presentation_path = os.path.join(DEMO_FILES_FOLDER, 'DEMO.pptx')
    
    if not os.path.isfile(presentation_path):
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Hello, World!"
        subtitle.text = "python-pptx was here!"
        prs.save(presentation_path)
    
    api_key = os.getenv("OPENAI_API_KEY")

    assert api_key is not None, "API key is not set in environment variables"

    slides_text = extract_text_from_presentation(presentation_path)
    summaries = await process_presentation(slides_text, api_key)
    text_to_json_file(summaries, presentation_path, output_dir=OUTPUTS_FOLDER)

    output_file = os.path.join(OUTPUTS_FOLDER, 'DEMO.json')
    assert os.path.isfile(output_file), "Output JSON file was not created"

    if os.path.isfile(output_file):
        os.remove(output_file)

@pytest.mark.asyncio
async def test_empty_presentation():
    presentation_path = os.path.join(DEMO_FILES_FOLDER, 'EMPTY.pptx')
    api_key = os.getenv("OPENAI_API_KEY")

    if not os.path.isfile(presentation_path):
        prs = Presentation()
        prs.save(presentation_path)

    assert api_key is not None, "API key is not set in environment variables"

    slides_text = extract_text_from_presentation(presentation_path)
    summaries = await process_presentation(slides_text, api_key)
    text_to_json_file(summaries, presentation_path, output_dir=OUTPUTS_FOLDER)
    
    output_file = os.path.join(OUTPUTS_FOLDER, 'EMPTY.json')
    assert os.path.isfile(output_file), "Output JSON file was not created"

    with open(output_file, 'r', encoding='utf-8') as file:
        data = json.load(file)
        assert len(data) == 0, "Output JSON should be empty for an empty presentation"

    if os.path.isfile(output_file):
        os.remove(output_file)
    if os.path.isfile(presentation_path):
        os.remove(presentation_path)

@pytest.mark.asyncio
async def test_missing_presentation():
    presentation_path = os.path.join(DEMO_FILES_FOLDER, 'MISSING.pptx')
    api_key = os.getenv("OPENAI_API_KEY")

    assert api_key is not None, "API key is not set in environment variables"

    with pytest.raises(FileNotFoundError):
        slides_text = extract_text_from_presentation(presentation_path)
        await process_presentation(slides_text, api_key)

@pytest.mark.asyncio
async def test_malformed_presentation():
    presentation_path = os.path.join(DEMO_FILES_FOLDER, 'MALFORMED.pptx')
    api_key = os.getenv("OPENAI_API_KEY")

    if not os.path.isfile(presentation_path):
        with open(presentation_path, 'w') as file:
            file.write("This is not a valid pptx file content")

    assert api_key is not None, "API key is not set in environment variables"

    with pytest.raises(Exception):
        slides_text = extract_text_from_presentation(presentation_path)
        await process_presentation(slides_text, api_key)

    if os.path.isfile(presentation_path):
        os.remove(presentation_path)

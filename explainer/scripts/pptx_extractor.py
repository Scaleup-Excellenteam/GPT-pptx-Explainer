from pptx import Presentation
import os
import json
import logging

logger = logging.getLogger('ExplainerLogger')
base_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '../../'))

def extract_slide_text(slide, counter: int) -> str:
    """
    Extracts text from a given slide.
    
    Parameters:
    - slide: The slide object to extract text from.
    - counter (int): The slide number.
    
    Returns:
    - str: The extracted text from the slide.
    """
    slide_text = f"Slide number: {counter}\n"
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            try:
                text = shape.text.strip()
                slide_text += text.encode('utf-8').decode('utf-8') + '\n'
            except UnicodeEncodeError:
                slide_text += "[Error: Unable to decode text]\n"
    return slide_text.strip()

def extract_text_from_presentation(presentation_path: str) -> list[tuple[int, str]]:
    """
    Extracts text from each slide in the presentation.

    Parameters:
    - presentation_path (str): The path to the presentation file.

    Returns:
    - list[tuple[int, str]]: A list of tuples where each tuple contains the slide number and its extracted text.
    
    Raises:
    - FileNotFoundError: If the presentation file does not exist.
    - Exception: If there is an error loading the presentation.
    """
    if not os.path.isfile(presentation_path):
        raise FileNotFoundError(f"The presentation file '{presentation_path}' does not exist.")
    
    try:
        presentation = Presentation(presentation_path)
    except Exception as e:
        raise Exception(f"Failed to load presentation: {e}")

    slides_text = []
    counter = 1
    for slide in presentation.slides:
        slide_text = extract_slide_text(slide, counter)
        if slide_text:
            slides_text.append((counter, slide_text))
        counter += 1
    logger.info(f"Extracted text from {len(slides_text)} slides in {presentation_path}")
    return slides_text

def text_to_json_file(summarized_text: dict[int, str], path: str, output_dir: str) -> None:
    """
    Saves summarized text to a JSON file.

    Parameters:
    - summarized_text (dict[int, str]): A dictionary where keys are slide numbers and values are summarized text.
    - path (str): The path to the original presentation file.
    - output_dir (str): The directory where the JSON file will be saved.
    """
    name = os.path.splitext(os.path.basename(path))[0]
    os.makedirs(output_dir, exist_ok=True)
    
    json_file_name = os.path.join(output_dir, name + '.json')
    with open(json_file_name, 'w', encoding='utf-8') as writer:
        json.dump(summarized_text, writer, indent=4, ensure_ascii=False)
    logger.info(f"Summarized text saved to {json_file_name}")

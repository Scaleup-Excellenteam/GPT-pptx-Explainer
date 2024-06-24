from pptx import Presentation
import os
import json

def extract_slide_text(slide, counter) -> str:
    """
    Extracts text from a slide and formats it with the slide number.

    Args:
        slide (Slide): The slide object from which to extract text.
        counter (int): The slide number.

    Returns:
        str: The extracted text formatted with the slide number.
    """
    slide_text = f"Slide number: {counter}\n"
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            try:
                text = shape.text.strip()
                slide_text += text.encode('utf-8').decode('utf-8') + '\n'
            except UnicodeEncodeError:
                slide_text += "[Error: Unable to decode text]\n"
    return slide_text.strip()

def extract_text_from_presentation(presentation_path: str) -> list[tuple[int, str]]:
    """
    Extracts text from all slides in a presentation.

    Args:
        presentation_path (str): The path to the presentation file.

    Returns:
        list[tuple[int, str]]: A list of tuples, each containing a slide number and its text content.
    
    Raises:
        FileNotFoundError: If the presentation file does not exist.
        Exception: If there is an error loading the presentation.
    """
    if not os.path.isfile(presentation_path):
        raise FileNotFoundError(f"The presentation file '{presentation_path}' does not exist.")
    
    try:
        presentation = Presentation(presentation_path)
    except Exception as e:
        raise Exception(f"Failed to load presentation: {e}")

    slides_text = []
    counter = 1
    for slide in presentation.slides:
        slide_text = extract_slide_text(slide, counter)
        if slide_text:
            slides_text.append((counter, slide_text))
        counter += 1
    print("Extracting text from slides...")
    return slides_text

def text_to_json_file(summarized_text: dict[int, str], path: str) -> None:
    """
    Saves the summarized text to a JSON file.

    Args:
        summarized_text (dict[int, str]): A dictionary containing slide numbers and their corresponding summaries.
        path (str): The path to the original presentation file.

    Returns:
        None
    """
    name = os.path.splitext(os.path.basename(path))[0]
    

    #print the name of the curr directory)
    print(os.path.basename(os.getcwd()))
    output_directory = os.path.join(os.getcwd(), 'outputs') 
    os.makedirs(output_directory, exist_ok=True)
    
    json_file_name = os.path.join(output_directory, name + '.json')
    with open(json_file_name, 'w', encoding='utf-8') as writer:
        json.dump(summarized_text, writer, indent=4, ensure_ascii=False)

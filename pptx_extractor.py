from pptx import Presentation
import os
import json
import asyncio
import ai_summery

async def fetch_summary(ai_key, slide_text):
    return await ai_summery.generate_summary(ai_key, slide_text)

def extract_slide_text(slide, counter):
    slide_text = 'Slide number: ' + str(counter) + '\n'
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            slide_text += shape.text.strip() + '\n'
    return slide_text.strip()

def create_tasks_from_presentation(presentation_path, ai_key):
    if not os.path.isfile(presentation_path):
        raise FileNotFoundError(f"The presentation file '{presentation_path}' does not exist.")
    
    try:
        presentation = Presentation(presentation_path)
    except Exception as e:
        raise Exception(f"Failed to load presentation: {e}")

    tasks = []
    counter = 1
    for slide in presentation.slides:
        slide_text = extract_slide_text(slide, counter)
        if slide_text:
            task = asyncio.create_task(fetch_summary(ai_key, slide_text))
            tasks.append((counter, task))
        counter += 1
    return tasks

async def extract_text_from_presentation(presentation_path, ai_key):
    tasks = create_tasks_from_presentation(presentation_path, ai_key)
    text = {}
    for counter, task in tasks:
        try:
            summary = await task
            text[counter] = summary
        except Exception as e:
            text[counter] = f"Error summarizing slide: {e}"
    return text

def text_to_json_file(summarized_text, path):
    name = os.path.splitext(os.path.basename(path))[0]
    output_directory = os.path.join(os.getcwd(), 'outputs')
    os.makedirs(output_directory, exist_ok=True)
    
    json_file_name = os.path.join(output_directory, name + '.json')
    with open(json_file_name, 'w') as writer:
        json.dump(summarized_text, writer, indent=4)

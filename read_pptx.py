import ai_summery
from pptx import Presentation
import asyncio
import json
from numpy import size

def read_presentation_text(presentation_path):
    presentation = Presentation(presentation_path)
    slides_text = {}
    counter = 1
    for slide in presentation.slides:
        slide_text = 'Slide number: ' + str(counter) + '\n'
        for shape in slide.shapes:
            if hasattr(shape, 'text'):  # check if shape has text attribute
                slide_text += shape.text.strip() + '\n'
        if slide_text.strip():  # check if slide_text is not empty
            slides_text[counter] = slide_text
            counter += 1
        # else ignore slide
    return slides_text

async def fetch_summary(ai_key, slide_text, timeout=10):
    return await ai_summery.generate_summary(ai_key, slide_text, timeout)

async def extract_text_from_presentation(presentation_path, ai_key, timeout=10):
    slides_text = read_presentation_text(presentation_path)
    tasks = []
    for counter, slide_text in slides_text.items():
        task = asyncio.create_task(fetch_summary(ai_key, slide_text, timeout))
        tasks.append((counter, task))
        
    text = {}
    for counter, task in tasks:
        summary = await task
        text[counter] = summary
    return text

def get_output_path(presentation_path):
    presentation_path = presentation_path.split('/') 
    output = presentation_path[size(presentation_path)-1] 
    output = output.split('.')
    return 'final-exercise-DSH93/outputs/' + output[0] + '.json'
    
def save_as_json(text, presentation_path):
    output_path = get_output_path(presentation_path)
    with open(output_path, 'w') as file:
        json.dump(text, file, indent=4)
        
if __name__ == '__main__':
    api_key = ai_summery.load_api_key('final-exercise-DSH93/openAi_key.txt')
    path = 'final-exercise-DSH93/example_pptx.pptx'
    text = asyncio.run(extract_text_from_presentation(path, api_key))
    save_as_json(text, path)
